import fitz            # pdf_parser
import trafilatura     # html_parser
import chardet         # txt_parser
import csv             # csv_parser
import json            # json_parser
import docx            # docx_parser
import pptx            # pptx_parser
import pandas as pd    # xlsx_parser
from pathlib import Path

def parse_pdf(pdf_path: str) -> list[dict]:
    blocks = []
    doc = fitz.open(pdf_path)

    for page_number, page in enumerate(doc, start=1):
        text = page.get_text().strip()
        if not text:
            continue  # skip empty/scanned pages

        blocks.append({
            "source_file": pdf_path,
            "block_type": "page",
            "text": text,
            "metadata": {"page_number": page_number},
        })

    doc.close()
    return blocks

def parse_html(html_path: str) -> list[dict]:
    with open(html_path, "r", encoding="utf-8", errors="ignore") as f:
        raw_html = f.read()

    text = trafilatura.extract(raw_html)
    if not text:
        return []

    return [{
        "source_file": html_path,
        "block_type": "page",
        "text": text.strip(),
        "metadata": {},
    }]

def parse_txt(txt_path: str) -> list[dict]:
    with open(txt_path, "rb") as f:
        raw = f.read()

    encoding = chardet.detect(raw)["encoding"] or "utf-8"
    text = raw.decode(encoding, errors="ignore").strip()

    if not text:
        return []

    return [{
        "source_file": txt_path,
        "block_type": "document",
        "text": text,
        "metadata": {},
    }]

def parse_csv(csv_path: str, rows_per_block) -> list[dict]:
    blocks = []

    with open(csv_path, "r", encoding="utf-8", errors="ignore", newline="") as f:
        reader = csv.DictReader(f)
        headers = reader.fieldnames or []
        batch = []

        for i, row in enumerate(reader):
            row_text = ", ".join(f"{k}: {v}" for k, v in row.items())
            batch.append(row_text)

            if len(batch) >= rows_per_block:
                blocks.append(_make_csv_block(csv_path, headers, batch, i))
                batch = []

        if batch:
            blocks.append(_make_csv_block(csv_path, headers, batch, "final"))

    return blocks
def _make_csv_block(source_file, headers, rows, batch_id) -> dict:
    return {
        "source_file": source_file,
        "block_type": "row_group",
        "text": "\n".join(rows),
        "metadata": {"columns": headers, "batch_id": batch_id},
    }

def parse_json(json_path: str) -> list[dict]:
    with open(json_path, "r", encoding="utf-8", errors="ignore") as f:
        raw = f.read().strip()

    records = _parse_records(raw)

    return [{
        "source_file": json_path,
        "block_type": "record",
        "text": json.dumps(record, ensure_ascii=False),
        "metadata": {"record_index": i},
    } for i, record in enumerate(records)]
def _parse_records(raw: str) -> list:
    try:
        data = json.loads(raw)
        return data if isinstance(data, list) else [data]
    except json.JSONDecodeError:
        # fall back to NDJSON (one record per line)
        return [json.loads(line) for line in raw.splitlines() if line.strip()]
    
def parse_docx(docx_path: str) -> list[dict]:
    doc = docx.Document(docx_path)
    blocks = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        is_heading = para.style.name.lower().startswith("heading")
        blocks.append({
            "source_file": docx_path,
            "block_type": "heading" if is_heading else "paragraph",
            "text": text,
            "metadata": {"style": para.style.name},
        })

    return blocks

def parse_pptx(pptx_path: str) -> list[dict]:
    prs = pptx.Presentation(pptx_path)
    blocks = []

    for i, slide in enumerate(prs.slides, start=1):
        texts = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if not texts:
            continue

        blocks.append({
            "source_file": pptx_path,
            "block_type": "slide",
            "text": "\n".join(texts),
            "metadata": {"slide_number": i},
        })

    return blocks

import pandas as pd


def parse_xlsx(xlsx_path: str, rows_per_block) -> list[dict]:
    blocks = []
    sheets = pd.read_excel(xlsx_path, sheet_name=None, dtype=str)

    for sheet_name, df in sheets.items():
        df = df.fillna("")
        rows = [
            ", ".join(f"{col}: {val}" for col, val in row.items())
            for _, row in df.iterrows()
        ]

        for i in range(0, len(rows), rows_per_block):
            batch = rows[i:i + rows_per_block]
            blocks.append({
                "source_file": xlsx_path,
                "block_type": "row_group",
                "text": "\n".join(batch),
                "metadata": {
                    "sheet": sheet_name,
                    "columns": list(df.columns),
                    "batch_id": i // rows_per_block,
                },
            })

    return blocks

Parsers = {
    ".pdf" : parse_pdf,
    ".html": parse_html,
    ".htm" : parse_html,
    ".txt" : parse_txt,
    ".csv" : parse_csv,
    ".json": parse_json,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".xlsx": parse_xlsx,
}

def parse_file(file_path: str, rows_per_block: int = 20):
    ext = Path(file_path).suffix.lower()
    parser = Parsers.get(ext)
    if parser is None:
        raise ValueError(f"No parser registered for extension: {ext}")
    return parser(file_path)